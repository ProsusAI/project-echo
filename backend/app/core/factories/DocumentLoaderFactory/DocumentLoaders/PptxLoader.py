import io
from pptx import Presentation

from app.core.factories.DocumentLoaderFactory.DocumentLoaders.AbstractDocumentLoader import (
    DocumentLoader,
)


class PptxLoader(DocumentLoader):
    file_types = ["ppt", "pptx"]
    mime_types = [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    def extract_text(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return " ".join(text)